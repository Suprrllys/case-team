#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Practicum Award 2026 — «Команда ИИ-агентов для решения бизнес-кейсов» (v3)
Builder под YAML-контракт designer-instruction.md (mode: regular).

Изменения v3:
- Брендбук Яндекс Практикум: жёлтый primary #FBCE07, чёрный hero, белый фон контента
- Шрифт YS Text / YS Display (fallback Inter / Manrope)
- 19 слайдов = 9 основных + 10 приложений A0..A9
- A0 «Источники» разбит на 3 подслайда (71 ссылка с кликабельными URL)
- Slide 8 «Команда» — карточка с placeholder-фото (инициалы СЛ если файла нет)
- Slide 9 «Финал» — формула differentiation Замесина
- 2 системы навигации: основная (slides 2-8) и appendix (slides 10-19)
- Pre-flight overflow check после save
"""

import os
import re
import zipfile
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
OUT_FILE = os.path.join(HERE, "practicum-award-2026-case-team.pptx")

# ── ПАЛИТРА YANDEX PRACTICUM ────────────────────────────────────────────
PRIMARY        = RGBColor(0xFB, 0xCE, 0x07)  # Yandex Yellow
PRIMARY_DARK   = RGBColor(0x00, 0x00, 0x00)  # Yandex Black
PRIMARY_SOFT   = RGBColor(0xFF, 0xF5, 0xCC)  # light yellow card fill
INK            = RGBColor(0x1A, 0x1A, 0x1A)  # main text
INK_2          = RGBColor(0x66, 0x66, 0x66)  # secondary text
BG             = RGBColor(0xFF, 0xFF, 0xFF)  # white slide background
ACCENT_POS     = RGBColor(0x00, 0xC8, 0x53)  # green metric
ACCENT_WARN    = RGBColor(0xFF, 0x6B, 0x00)  # orange warn
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
DIVIDER        = RGBColor(0xE5, 0xE5, 0xE5)
GREY_PALE      = RGBColor(0xF3, 0xF3, 0xF3)
LINK_BLUE      = RGBColor(0x0B, 0x57, 0xD0)  # hyperlink

FONT_DISPLAY = "YS Display"   # fallback: Inter / Manrope (system-dependent)
FONT_BODY    = "YS Text"      # fallback: Inter / Manrope
FONT_MONO    = "JetBrains Mono"

# ── КАНВА 16:9 ──────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── НАВИГАЦИЯ (основная: slides 2..8, табы — 6 разделов) ────────────────
MAIN_NAV_TABS = ["Резюме", "Анализ", "Концепция", "Реализация", "Эффект", "Команда"]
# Slide → active tab index (None = no nav)
SLIDE_TO_TAB = {
    2: 0,  # Exec Summary → Резюме
    3: 1,  # Анализ
    # 4: hero (без табов)
    5: 3,  # Реализация
    6: 4,  # Эффект (экономика)
    7: 4,  # Эффект (что меняется)
    8: 5,  # Команда
}


# ── ХЕЛПЕРЫ ─────────────────────────────────────────────────────────────

def add_run(paragraph, text, *, size=14, bold=False, color=INK, font=FONT_BODY,
            italic=False, hyperlink=None):
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    # Явная фиксация шрифта через rPr (чтобы не подменился на Calibri)
    try:
        rPr = run._r.get_or_add_rPr()
        # latin
        latin = rPr.find(qn("a:latin"))
        if latin is None:
            latin = etree.SubElement(rPr, qn("a:latin"))
        latin.set("typeface", font)
        # east-asian
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", font)
    except Exception:
        pass
    if hyperlink:
        run.hyperlink.address = hyperlink
    return run


def set_text(tf, text, *, size=14, bold=False, color=INK, font=FONT_BODY,
             align=PP_ALIGN.LEFT, line_spacing=1.15, italic=False):
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for r in list(p.runs):
        r._r.getparent().remove(r._r)
    add_runs_with_links(p, text, size=size, bold=bold, color=color, font=font,
                        italic=italic)
    return p


def add_paragraph(tf, text, *, size=14, bold=False, color=INK, font=FONT_BODY,
                  align=PP_ALIGN.LEFT, line_spacing=1.15, bullet=False,
                  space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line_spacing
    if space_before:
        p.space_before = Pt(space_before)
    if bullet:
        add_run(p, "•  ", size=size, color=PRIMARY_DARK, font=font, bold=True)
    add_runs_with_links(p, text, size=size, bold=bold, color=color, font=font,
                        italic=italic)
    return p


MD_LINK_RE = re.compile(r'\[([^\]]+)\]\(([^)]+)\)')


def add_runs_with_links(paragraph, text, *, size=14, bold=False, color=INK,
                        font=FONT_BODY, italic=False):
    """Парсит markdown-ссылки [текст](url) в строке и рендерит их как
    кликабельные runs с синим цветом + подчёркиванием. Остальной текст —
    обычный run с заданным color.

    Возвращает список созданных runs.
    """
    runs = []
    last_end = 0
    for m in MD_LINK_RE.finditer(text):
        # Текст до ссылки
        if m.start() > last_end:
            plain = text[last_end:m.start()]
            if plain:
                runs.append(add_run(paragraph, plain, size=size, bold=bold,
                                    color=color, font=font, italic=italic))
        # Ссылка
        link_text = m.group(1)
        link_url = m.group(2)
        r = add_run(paragraph, link_text, size=size, bold=bold,
                    color=LINK_BLUE, font=font, italic=italic,
                    hyperlink=link_url)
        # Подчёркивание для визуальной идентификации ссылки
        try:
            r.font.underline = True
        except Exception:
            pass
        runs.append(r)
        last_end = m.end()
    # Остаток текста после последней ссылки
    if last_end < len(text):
        tail = text[last_end:]
        if tail:
            runs.append(add_run(paragraph, tail, size=size, bold=bold,
                                color=color, font=font, italic=italic))
    # Если в строке не было ни одной ссылки — добавить как обычный текст
    if not runs:
        runs.append(add_run(paragraph, text, size=size, bold=bold,
                            color=color, font=font, italic=italic))
    return runs


def add_rect(slide, x, y, w, h, *, fill=None, line=None, radius=None,
             line_width=0.75):
    if radius is not None:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    # Удалить тени
    try:
        spPr = shape._element.spPr
        for el in spPr.findall(qn("a:effectLst")):
            spPr.remove(el)
        etree.SubElement(spPr, qn("a:effectLst"))
    except Exception:
        pass
    return shape


def add_oval(slide, x, y, w, h, *, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    try:
        spPr = shape._element.spPr
        for el in spPr.findall(qn("a:effectLst")):
            spPr.remove(el)
        etree.SubElement(spPr, qn("a:effectLst"))
    except Exception:
        pass
    return shape


def add_textbox(slide, x, y, w, h, *, text="", size=14, bold=False, color=INK,
                font=FONT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                line_spacing=1.15, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    set_text(tf, text, size=size, bold=bold, color=color, font=font,
             align=align, line_spacing=line_spacing, italic=italic)
    return tb


def paint_background(slide, color=BG):
    bg = add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=color)
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)
    return bg


def draw_main_nav(slide, active_idx):
    """Главная навигация на slides 2-8 (кроме cover и final)."""
    tab_y = Inches(0.30)
    tab_h = Inches(0.36)
    margin_x = Inches(0.55)
    avail_w = SLIDE_W - margin_x * 2
    n = len(MAIN_NAV_TABS)
    gap = Inches(0.08)
    tab_w = (avail_w - gap * (n - 1)) / n
    x = margin_x
    for i, label in enumerate(MAIN_NAV_TABS):
        if i == active_idx:
            add_rect(slide, x, tab_y, tab_w, tab_h, fill=PRIMARY, radius=0.5)
            color = PRIMARY_DARK
            bold = True
        else:
            add_rect(slide, x, tab_y, tab_w, tab_h, fill=GREY_PALE, radius=0.5)
            color = INK_2
            bold = False
        tb = add_textbox(slide, x, tab_y, tab_w, tab_h, text=label, size=11,
                         bold=bold, color=color, align=PP_ALIGN.CENTER,
                         anchor=MSO_ANCHOR.MIDDLE)
        x += tab_w + gap


def draw_appendix_nav(slide, index, total):
    """Навигация для приложений: 'Приложение A{N} из A{TOTAL}'."""
    tab_y = Inches(0.30)
    tab_h = Inches(0.36)
    label = f"Приложение A{index} из A{total}"
    # Левый бейдж
    badge_w = Inches(3.6)
    add_rect(slide, Inches(0.55), tab_y, badge_w, tab_h, fill=PRIMARY_DARK,
             radius=0.5)
    add_textbox(slide, Inches(0.55), tab_y, badge_w, tab_h, text=label,
                size=11, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    # Маркер «приложения» справа
    right_label = "Q&A · опционально по запросу жюри"
    add_textbox(slide, SLIDE_W - Inches(5.05), tab_y, Inches(4.5), tab_h,
                text=right_label, size=10, color=INK_2, align=PP_ALIGN.RIGHT,
                anchor=MSO_ANCHOR.MIDDLE)


def draw_header(slide, slide_num, title, subtitle, *, dark=False):
    """Заголовок-тезис + подзаголовок-вывод серым."""
    title_y = Inches(0.85)
    title_h = Inches(1.20)
    title_color = WHITE if dark else INK
    sub_color = RGBColor(0xC9, 0xC9, 0xC9) if dark else INK_2
    add_textbox(slide, Inches(0.55), title_y, SLIDE_W - Inches(1.10), title_h,
                text=title, size=24, bold=True, color=title_color,
                font=FONT_DISPLAY, line_spacing=1.10)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(2.10),
                    SLIDE_W - Inches(1.10), Inches(0.70),
                    text=subtitle, size=14, color=sub_color,
                    line_spacing=1.25)
    # Тонкая линия-разделитель
    if not dark:
        add_rect(slide, Inches(0.55), Inches(2.78), Inches(0.60), Pt(2),
                 fill=PRIMARY)


def draw_footer(slide, slide_num, total=21, *, dark=False):
    color = RGBColor(0x9A, 0x9A, 0x9A) if dark else INK_2
    add_textbox(slide, Inches(0.55), Inches(7.10), Inches(6.0), Inches(0.30),
                text="Practicum Award 2026 · Сергей Леонтьев",
                size=9, color=color)
    add_textbox(slide, SLIDE_W - Inches(2.55), Inches(7.10),
                Inches(2.0), Inches(0.30),
                text=f"{slide_num} / {total}",
                size=9, color=color, align=PP_ALIGN.RIGHT)


def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ────────────────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ────────────────────────────────────────────────────────────────────────

def build_slide_01_cover(prs):
    """Cover — чёрный фон, жёлтый акцент, белая типографика."""
    s = add_blank(prs)
    paint_background(s, color=PRIMARY_DARK)

    # Жёлтая полоска-акцент слева сверху
    add_rect(s, Inches(0.55), Inches(0.55), Inches(0.65), Inches(0.18),
             fill=PRIMARY)

    # Метка номинации
    add_textbox(s, Inches(0.55), Inches(0.85), Inches(9.0), Inches(0.4),
                text="PRACTICUM AWARD 2026  ·  ТЕХНО-ПРОДУКТОВЫЙ ВКЛАД",
                size=11, bold=True, color=PRIMARY, font=FONT_DISPLAY)

    # Главный заголовок
    add_textbox(s, Inches(0.55), Inches(1.85), Inches(12.3), Inches(2.20),
                text="Альтернатива стратегическому консалтингу через мультиагентную AI-команду",
                size=44, bold=True, color=WHITE, font=FONT_DISPLAY,
                line_spacing=1.10)

    # Подзаголовок — главный тезис
    add_textbox(s, Inches(0.55), Inches(4.20), Inches(12.3), Inches(1.0),
                text="Полный пакет материалов под защиту — за 4 часа, в стоимости подписки на нейросеть.",
                size=20, color=PRIMARY, line_spacing=1.30)

    # Три факта-плашки внизу
    cards = [
        ("4 часа", "вместо 5-6 дней"),
        ("11 тыс ₽", "vs 99 тыс ₽ in-house"),
        ("84-98%", "экономии (sensitivity ±20%)"),
    ]
    card_w = Inches(3.95)
    card_h = Inches(1.20)
    gap = Inches(0.20)
    total_w = card_w * 3 + gap * 2
    x0 = (SLIDE_W - total_w) / 2
    for i, (big, small) in enumerate(cards):
        x = x0 + i * (card_w + gap)
        y = Inches(5.55)
        add_rect(s, x, y, card_w, card_h, fill=RGBColor(0x1A, 0x1A, 0x1A),
                 radius=0.10)
        add_textbox(s, x, y + Inches(0.15), card_w, Inches(0.62),
                    text=big, size=32, bold=True, color=PRIMARY,
                    font=FONT_DISPLAY, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, x, y + Inches(0.75), card_w, Inches(0.40),
                    text=small, size=12, color=WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # Автор внизу
    add_textbox(s, Inches(0.55), Inches(7.00), Inches(12.0), Inches(0.35),
                text="Сергей Леонтьев  ·  Product Manager, Нинтегра  ·  МАИ 2025",
                size=12, color=RGBColor(0xC9, 0xC9, 0xC9))


def build_slide_02_exec(prs):
    """Exec Summary — 4 блока в 2 колонки."""
    s = add_blank(prs)
    paint_background(s)
    draw_main_nav(s, SLIDE_TO_TAB[2])
    draw_header(s, 2,
                "Готовим полный пакет материалов решения бизнес-кейса под защиту — за стоимость подписки на нейросеть",
                "За 4 часа вместо 5-6 дней, без дополнительной оплаты для пользователей Claude, с воспроизводимой методологией.")

    # 4 блока 2×2
    blocks = [
        ("АНАЛИЗ", "266 тыс компаний МСП в РФ принимают стратегические решения 2-6 раз в год; 5 текущих способов не дают сшитого пакета под защиту."),
        ("КОНЦЕПЦИЯ", "5 Claude-агентов готовят презентации, финмодели, аналитику и позиционирование. Решение go/no-go и защита перед людьми остаются за человеком."),
        ("РЕАЛИЗАЦИЯ", "5 ролей · обучены на решениях для Axenix (Cup Moscow 2024) и Askona (Inno Case Hack 2025) · 26 методологических логик — [открытый репозиторий](https://github.com/Suprrllys/case-team)."),
        ("ЭФФЕКТ", "4 часа вместо 5-6 дней · в стоимости подписки на нейросеть · устойчиво 84-98% экономии при ±20% к параметрам."),
    ]
    card_w = Inches(6.05)
    card_h = Inches(1.85)
    gap = Inches(0.20)
    x_left = Inches(0.55)
    x_right = x_left + card_w + gap
    y_top = Inches(3.15)
    y_bot = y_top + card_h + gap
    coords = [(x_left, y_top), (x_right, y_top), (x_left, y_bot), (x_right, y_bot)]

    for (label, body), (x, y) in zip(blocks, coords):
        add_rect(s, x, y, card_w, card_h, fill=PRIMARY_SOFT, radius=0.06)
        # Жёлтая полоска
        add_rect(s, x, y, Inches(0.10), card_h, fill=PRIMARY)
        # Caption
        add_textbox(s, x + Inches(0.30), y + Inches(0.18),
                    card_w - Inches(0.45), Inches(0.32),
                    text=label, size=10, bold=True, color=PRIMARY_DARK,
                    font=FONT_DISPLAY)
        # Body
        add_textbox(s, x + Inches(0.30), y + Inches(0.55),
                    card_w - Inches(0.45), card_h - Inches(0.70),
                    text=body, size=13, color=INK, line_spacing=1.30)

    draw_footer(s, 2)


def build_slide_03_problem(prs):
    """Анализ — рынок, ситуации, конкуренты."""
    s = add_blank(prs)
    paint_background(s)
    draw_main_nav(s, SLIDE_TO_TAB[3])
    draw_header(s, 3,
                "266 тыс компаний МСП в РФ принимают стратегические решения 2-6 раз в год — и каждый раз должны защитить их перед инвестором, советом или партнёром",
                "Из 9 проанализированных альтернатив ни одна не даёт сшитого пакета (research + стратегия + финмодель + презентация) под одной методологической логикой.")

    # Левая колонка — ключевые факты (numbers + контекст)
    x_left = Inches(0.55)
    y_blk = Inches(3.10)
    left_w = Inches(4.40)

    # Три числовых якоря
    metrics = [
        ("266 тыс", "компаний МСП в РФ"),
        ("2-6 раз / год", "принимают крупное стратегическое решение"),
        ("18-80 млрд ₽", "недополученного спроса в год"),
    ]
    for i, (num, ctx) in enumerate(metrics):
        y = y_blk + i * Inches(1.25)
        add_rect(s, x_left, y, left_w, Inches(1.10), fill=PRIMARY_SOFT,
                 radius=0.06)
        add_textbox(s, x_left + Inches(0.25), y + Inches(0.12), left_w - Inches(0.50),
                    Inches(0.55), text=num, size=26, bold=True,
                    color=PRIMARY_DARK, font=FONT_DISPLAY)
        add_textbox(s, x_left + Inches(0.25), y + Inches(0.66), left_w - Inches(0.50),
                    Inches(0.40), text=ctx, size=12, color=INK_2,
                    line_spacing=1.25)

    # Правая колонка — таблица альтернатив
    x_right = Inches(5.20)
    right_w = Inches(7.60)
    y0 = Inches(3.10)
    # Заголовок таблицы
    add_textbox(s, x_right, y0, right_w, Inches(0.32),
                text="ЧТО КЛИЕНТ НАНИМАЕТ СЕЙЧАС  ·  ПОЧЕМУ НЕ ПОДХОДИТ",
                size=10, bold=True, color=PRIMARY_DARK, font=FONT_DISPLAY)

    rows = [
        ("Делать самому", "5-6 дней", "нет компетенций в 5 ролях"),
        ("In-house команда из 5", "99k + 700k/мес ФОТ", "неподъёмно для МСП"),
        ("Big4-tier агентство", "350k — 1.5M ₽", "цена за репутацию партнёра"),
        ("Фрилансеры с биржи", "50-150k ₽", "никто не сшивает в историю"),
        ("ChatGPT / один AI-чат", "подписка", "текст, не пакет; нет формул"),
    ]
    row_h = Inches(0.55)
    col_w = [Inches(2.80), Inches(1.80), Inches(3.00)]
    cur_y = y0 + Inches(0.42)
    # Заголовки колонок
    headers = ["Альтернатива", "Стоимость", "Главный минус"]
    cx = x_right
    for i, hdr in enumerate(headers):
        add_textbox(s, cx, cur_y, col_w[i], Inches(0.30),
                    text=hdr, size=9, bold=True, color=INK_2,
                    font=FONT_DISPLAY)
        cx += col_w[i]
    cur_y += Inches(0.35)
    add_rect(s, x_right, cur_y, right_w, Pt(1), fill=DIVIDER)
    cur_y += Inches(0.05)

    for row in rows:
        cx = x_right
        for i, cell in enumerate(row):
            bold = (i == 0)
            color = INK if i == 0 else INK_2
            size = 11 if i == 0 else 10
            add_textbox(s, cx, cur_y + Inches(0.06), col_w[i], row_h,
                        text=cell, size=size, bold=bold, color=color,
                        line_spacing=1.20)
            cx += col_w[i]
        cur_y += row_h
        add_rect(s, x_right, cur_y, right_w, Pt(0.5), fill=DIVIDER)

    draw_footer(s, 3)


def build_slide_04_concept(prs):
    """Концепция — hero accent:dark."""
    s = add_blank(prs)
    paint_background(s, color=PRIMARY_DARK)

    # Тонкая жёлтая полоска вверху
    add_rect(s, Inches(0.55), Inches(0.55), Inches(0.65), Inches(0.16),
             fill=PRIMARY)
    add_textbox(s, Inches(0.55), Inches(0.80), Inches(8.0), Inches(0.35),
                text="КОНЦЕПЦИЯ  ·  АРХИТЕКТУРНЫЕ РЕШЕНИЯ",
                size=11, bold=True, color=PRIMARY, font=FONT_DISPLAY)

    # Заголовок-тезис
    add_textbox(s, Inches(0.55), Inches(1.30), Inches(12.3), Inches(1.45),
                text="5 Claude-агентов готовят полный пакет — а решение и защита остаются за человеком",
                size=28, bold=True, color=WHITE, font=FONT_DISPLAY,
                line_spacing=1.12)

    add_textbox(s, Inches(0.55), Inches(2.80), Inches(12.3), Inches(0.55),
                text="Куратор смысла внутри команды + физические артефакты на выходе + дотренировка на двух реальных кейс-чемпионатах — этого нет ни у одного из 9 проанализированных конкурентов.",
                size=13, color=RGBColor(0xCC, 0xCC, 0xCC), line_spacing=1.30)

    # 2 колонки: что делаем / что остаётся за человеком
    col_w = Inches(6.05)
    col_h = Inches(1.85)
    gap = Inches(0.20)
    y_cols = Inches(3.55)

    # Левая — что делаем
    add_rect(s, Inches(0.55), y_cols, col_w, col_h,
             fill=RGBColor(0x1A, 0x1A, 0x1A), radius=0.06)
    add_textbox(s, Inches(0.75), y_cols + Inches(0.15), col_w - Inches(0.40),
                Inches(0.32), text="ЧТО МЫ ДЕЛАЕМ ПОЛНОСТЬЮ",
                size=10, bold=True, color=PRIMARY, font=FONT_DISPLAY)
    tb = add_textbox(s, Inches(0.75), y_cols + Inches(0.50),
                     col_w - Inches(0.40), col_h - Inches(0.60),
                     text="", size=11, color=WHITE, line_spacing=1.30)
    tf = tb.text_frame
    items_left = [
        "Презентация .pptx — с навигацией, готова к показу",
        "Финмодель .xlsx — с формулами и sensitivity",
        "Аналитика — рынок, конкуренты, ЦА с источниками",
        "Концепция и позиционирование",
        "Decision Log — обоснование каждой развилки",
    ]
    set_text(tf, "•  " + items_left[0], size=11, color=WHITE, line_spacing=1.30)
    for it in items_left[1:]:
        add_paragraph(tf, "•  " + it, size=11, color=WHITE, line_spacing=1.30,
                      space_before=2)

    # Правая — что за человеком
    x_r = Inches(0.55) + col_w + gap
    add_rect(s, x_r, y_cols, col_w, col_h,
             fill=RGBColor(0x1A, 0x1A, 0x1A), radius=0.06)
    add_textbox(s, x_r + Inches(0.20), y_cols + Inches(0.15),
                col_w - Inches(0.40), Inches(0.32),
                text="ЧТО ОСТАЁТСЯ ЗА ЧЕЛОВЕКОМ",
                size=10, bold=True, color=PRIMARY, font=FONT_DISPLAY)
    tb = add_textbox(s, x_r + Inches(0.20), y_cols + Inches(0.50),
                     col_w - Inches(0.40), col_h - Inches(0.60),
                     text="", size=11, color=WHITE, line_spacing=1.30)
    tf = tb.text_frame
    items_right = [
        "Финальное go/no-go решение",
        "Защита перед инвестором / советом — живая коммуникация",
        "Операционное внедрение стратегии",
        "Договорённости с конкретными партнёрами",
    ]
    set_text(tf, "•  " + items_right[0], size=11, color=WHITE, line_spacing=1.30)
    for it in items_right[1:]:
        add_paragraph(tf, "•  " + it, size=11, color=WHITE, line_spacing=1.30,
                      space_before=2)

    # 4 архитектурных решения — 4 узких карточки внизу
    arch = [
        ("01", "Куратор смысла", "Менеджер-агент решает, что класть на слайды. Это ответ на 85% failure rate автономных агентов (Devin)."),
        ("02", "Физические артефакты", ".xlsx с формулами, .pptx готов к показу, .md Decision Log — а не текст в чате."),
        ("03", "Дотренировка на 2 кейс-чемпионатах", "Решения для Axenix (Cup Moscow 2024) и Askona (Inno Case Hack 2025) → 26 методологических логик в [открытом репозитории](https://github.com/Suprrllys/case-team)."),
        ("04", "Иерархия источников", "Эталон → бизнес-классика (Osterwalder, MoSCoW, AS-IS/TO-BE) → эвристика. Каждое правило ссылается на источник."),
    ]
    aw = Inches(3.05)
    ah = Inches(1.40)
    ag = Inches(0.10)
    ax0 = (SLIDE_W - (aw * 4 + ag * 3)) / 2
    ay = Inches(5.55)
    for i, (num, ttl, body) in enumerate(arch):
        x = ax0 + i * (aw + ag)
        add_rect(s, x, ay, aw, ah, fill=RGBColor(0x2A, 0x2A, 0x2A),
                 radius=0.08)
        add_textbox(s, x + Inches(0.18), ay + Inches(0.10), Inches(0.50),
                    Inches(0.30), text=num, size=14, bold=True, color=PRIMARY,
                    font=FONT_DISPLAY)
        add_textbox(s, x + Inches(0.18), ay + Inches(0.38), aw - Inches(0.36),
                    Inches(0.32), text=ttl, size=11, bold=True, color=WHITE,
                    font=FONT_DISPLAY)
        add_textbox(s, x + Inches(0.18), ay + Inches(0.70), aw - Inches(0.36),
                    ah - Inches(0.78), text=body, size=9, color=RGBColor(0xBB, 0xBB, 0xBB),
                    line_spacing=1.25)

    draw_footer(s, 4, dark=True)


def build_slide_05_realization(prs):
    """Реализация — схема + 5 ролей."""
    s = add_blank(prs)
    paint_background(s)
    draw_main_nav(s, SLIDE_TO_TAB[5])
    draw_header(s, 5,
                "5 Claude-агентов с дисциплиной коммуникации через файлы, обученных на решениях для Axenix и Askona",
                "Архитектура устроена так, что специалисты не пишут друг другу — всё проходит через менеджера. Это снимает риск рассинхрона между маркетингом, финансами и аналитикой.")

    # Схема — 5 узлов: case-manager наверху, под ним 3 specialist'а, под ними designer
    y_top = Inches(3.10)
    mgr_w = Inches(3.20)
    mgr_h = Inches(0.85)
    mgr_x = (SLIDE_W - mgr_w) / 2
    add_rect(s, mgr_x, y_top, mgr_w, mgr_h, fill=PRIMARY, radius=0.10)
    add_textbox(s, mgr_x, y_top + Inches(0.10), mgr_w, Inches(0.32),
                text="case-manager",
                size=13, bold=True, color=PRIMARY_DARK, font=FONT_DISPLAY,
                align=PP_ALIGN.CENTER)
    add_textbox(s, mgr_x, y_top + Inches(0.45), mgr_w, Inches(0.32),
                text="куратор смысла · PM-оркестратор",
                size=10, color=PRIMARY_DARK, align=PP_ALIGN.CENTER)

    # 3 специалиста под менеджером
    spec_y = Inches(4.30)
    spec_w = Inches(2.60)
    spec_h = Inches(0.85)
    spec_gap = Inches(0.30)
    total_spec_w = spec_w * 3 + spec_gap * 2
    spec_x0 = (SLIDE_W - total_spec_w) / 2
    specs = [
        ("case-analyst", "research"),
        ("case-financier", "финмодель"),
        ("case-marketer", "AJTBD"),
    ]
    for i, (name, sub) in enumerate(specs):
        x = spec_x0 + i * (spec_w + spec_gap)
        add_rect(s, x, spec_y, spec_w, spec_h, fill=PRIMARY_SOFT,
                 line=PRIMARY_DARK, line_width=0.5, radius=0.10)
        add_textbox(s, x, spec_y + Inches(0.12), spec_w, Inches(0.32),
                    text=name, size=12, bold=True, color=PRIMARY_DARK,
                    font=FONT_DISPLAY, align=PP_ALIGN.CENTER)
        add_textbox(s, x, spec_y + Inches(0.47), spec_w, Inches(0.30),
                    text=sub, size=10, color=INK_2, align=PP_ALIGN.CENTER)

    # Designer под специалистами
    des_y = Inches(5.50)
    des_x = mgr_x
    add_rect(s, des_x, des_y, mgr_w, mgr_h, fill=PRIMARY_DARK, radius=0.10)
    add_textbox(s, des_x, des_y + Inches(0.10), mgr_w, Inches(0.32),
                text="case-designer", size=13, bold=True, color=PRIMARY,
                font=FONT_DISPLAY, align=PP_ALIGN.CENTER)
    add_textbox(s, des_x, des_y + Inches(0.45), mgr_w, Inches(0.32),
                text="конвертер YAML → .pptx",
                size=10, color=WHITE, align=PP_ALIGN.CENTER)

    # Линии-связи: manager → каждый specialist, каждый specialist → designer
    mgr_bot_x = mgr_x + mgr_w / 2
    mgr_bot_y = y_top + mgr_h
    des_top_y = des_y
    for i in range(3):
        sx = spec_x0 + i * (spec_w + spec_gap) + spec_w / 2
        # manager → specialist
        ln = s.shapes.add_connector(1, mgr_bot_x, mgr_bot_y, sx, spec_y)
        ln.line.color.rgb = INK_2
        ln.line.width = Pt(1.0)
        # specialist → designer
        ln2 = s.shapes.add_connector(1, sx, spec_y + spec_h, mgr_bot_x, des_top_y)
        ln2.line.color.rgb = DIVIDER
        ln2.line.width = Pt(0.5)

    # Сноска про дотренировку
    note_y = Inches(6.55)
    add_textbox(s, Inches(0.55), note_y, Inches(12.2), Inches(0.50),
                text="Дотренировка — на решениях для Axenix (Cup Moscow 2024) и Askona (Inno Case Hack 2025) с кейс-чемпионатов. Извлечено 26 методологических логик — материалы в [открытом репозитории github.com/Suprrllys/case-team](https://github.com/Suprrllys/case-team).",
                size=11, color=INK_2, line_spacing=1.30)

    draw_footer(s, 5)


def build_slide_06_economics(prs):
    """Эффект — экономика: бары + sensitivity."""
    s = add_blank(prs)
    paint_background(s)
    draw_main_nav(s, SLIDE_TO_TAB[6])
    draw_header(s, 6,
                "Для пользователя Claude — без дополнительной оплаты; внутри команды 1 кейс стоит 11 тыс ₽ при 99 тыс ₽ in-house и 396 тыс ₽ через агентство",
                "Категориальный сдвиг — стратегический пакет встроен в стоимость подписки на нейросеть, которая уже оплачена клиентом для других задач.")

    # Левая колонка — главный факт-якорь крупно
    x_l = Inches(0.55)
    y_l = Inches(3.10)
    l_w = Inches(5.20)
    add_rect(s, x_l, y_l, l_w, Inches(2.10), fill=PRIMARY, radius=0.08)
    add_textbox(s, x_l + Inches(0.25), y_l + Inches(0.15), l_w - Inches(0.50),
                Inches(0.40), text="ДЛЯ ПОЛЬЗОВАТЕЛЯ CLAUDE",
                size=11, bold=True, color=PRIMARY_DARK, font=FONT_DISPLAY)
    add_textbox(s, x_l + Inches(0.25), y_l + Inches(0.55), l_w - Inches(0.50),
                Inches(1.0), text="0 ₽", size=60, bold=True,
                color=PRIMARY_DARK, font=FONT_DISPLAY)
    add_textbox(s, x_l + Inches(0.25), y_l + Inches(1.55), l_w - Inches(0.50),
                Inches(0.45), text="дополнительного расхода — подписка уже оплачена для других задач",
                size=11, color=PRIMARY_DARK, line_spacing=1.25)

    # Под ним — 4 часа vs 5-6 дней
    add_rect(s, x_l, y_l + Inches(2.30), l_w, Inches(1.10),
             fill=PRIMARY_SOFT, radius=0.06)
    add_textbox(s, x_l + Inches(0.25), y_l + Inches(2.42), l_w - Inches(0.50),
                Inches(0.35), text="ВРЕМЯ", size=10, bold=True,
                color=PRIMARY_DARK, font=FONT_DISPLAY)
    add_textbox(s, x_l + Inches(0.25), y_l + Inches(2.72), l_w - Inches(0.50),
                Inches(0.45), text="4 часа", size=28, bold=True, color=PRIMARY_DARK,
                font=FONT_DISPLAY)
    add_textbox(s, x_l + Inches(2.40), y_l + Inches(2.85), l_w - Inches(2.60),
                Inches(0.45), text="вместо 5-6 дней команды",
                size=11, color=INK_2)

    # Правая колонка — бар-чарт стоимости 1 кейса
    x_r = Inches(6.10)
    r_w = Inches(6.70)
    add_textbox(s, x_r, Inches(3.10), r_w, Inches(0.32),
                text="СТОИМОСТЬ 1 КЕЙСА (ВНУТРЕННЯЯ ЭКОНОМИКА, ₽)",
                size=10, bold=True, color=INK_2, font=FONT_DISPLAY)

    bars = [
        ("Агенты (менеджер + tokens)", 10930, ACCENT_POS),
        ("In-house команда из 5", 99000, INK_2),
        ("Big4-tier агентство", 396000, RGBColor(0xC0, 0x39, 0x2B)),
    ]
    max_val = max(v for _, v, _ in bars)
    bar_y0 = Inches(3.55)
    bar_h = Inches(0.55)
    bar_gap = Inches(0.32)
    label_w = Inches(2.80)
    bar_area_w = r_w - label_w - Inches(0.20)
    for i, (lab, val, col) in enumerate(bars):
        y = bar_y0 + i * (bar_h + bar_gap)
        # Лейбл
        add_textbox(s, x_r, y + Inches(0.05), label_w, bar_h - Inches(0.10),
                    text=lab, size=11, color=INK, line_spacing=1.20)
        # Бар
        bx = x_r + label_w
        bw = int(bar_area_w * (val / max_val))
        add_rect(s, bx, y + Inches(0.08), bw, bar_h - Inches(0.16),
                 fill=col, radius=0.10)
        # Значение справа от бара
        add_textbox(s, bx + bw + Inches(0.10), y + Inches(0.10),
                    Inches(1.5), bar_h - Inches(0.20),
                    text=f"{val:,} ₽".replace(",", " "),
                    size=12, bold=True, color=INK, font=FONT_DISPLAY)

    # Sensitivity внизу
    sens_y = Inches(6.05)
    add_textbox(s, x_r, sens_y, r_w, Inches(0.32),
                text="SENSITIVITY ±20% — экономия vs in-house",
                size=10, bold=True, color=INK_2, font=FONT_DISPLAY)
    sens = [("Базовый", "89%"), ("Junior in-house", "84%"),
            ("Big4 sourcing", "97-98%"), ("+20% токенов", "88%")]
    sw = (r_w - Inches(0.20) * 3) / 4
    sy = sens_y + Inches(0.40)
    for i, (lab, val) in enumerate(sens):
        x = x_r + i * (sw + Inches(0.20))
        add_rect(s, x, sy, sw, Inches(0.75), fill=PRIMARY_SOFT, radius=0.10)
        add_textbox(s, x, sy + Inches(0.06), sw, Inches(0.32),
                    text=val, size=18, bold=True, color=ACCENT_POS,
                    font=FONT_DISPLAY, align=PP_ALIGN.CENTER)
        add_textbox(s, x, sy + Inches(0.42), sw, Inches(0.30),
                    text=lab, size=9, color=INK_2, align=PP_ALIGN.CENTER)

    draw_footer(s, 6)


def build_slide_07_changes(prs):
    """Что меняется + низкий порог входа (две колонки)."""
    s = add_blank(prs)
    paint_background(s)
    draw_main_nav(s, SLIDE_TO_TAB[7])
    draw_header(s, 7,
                "Что меняется относительно существующих способов и что снижает порог входа",
                "Главный страх клиента — не «дорого», а «обнаружится фейк под вопросами стейкхолдера». Поэтому каждая цифра защищена источником.")

    # 2 колонки таблиц
    y0 = Inches(3.10)
    col_w = Inches(6.10)
    col_gap = Inches(0.20)
    x_l = Inches(0.55)
    x_r = x_l + col_w + col_gap

    # Левая — Что меняется
    add_textbox(s, x_l, y0, col_w, Inches(0.36),
                text="ЧТО МЕНЯЕТСЯ ОТНОСИТЕЛЬНО СУЩЕСТВУЮЩИХ СПОСОБОВ",
                size=10, bold=True, color=PRIMARY_DARK, font=FONT_DISPLAY)
    left_rows = [
        ("Делать самому", "4 часа реального времени; качество в 5 ролях"),
        ("In-house команда из 5", "Не нужен 1-2-месячный найм; пакет за стоимость подписки"),
        ("Big4 (350k-1.5M ₽)", "4 часа вместо недель; прозрачная методология"),
        ("Фрилансеры с биржи", "Куратор смысла внутри → единая методология"),
        ("ChatGPT / один AI-чат", ".xlsx и .pptx с формулами, не текст; источник на каждую цифру"),
    ]
    row_y = y0 + Inches(0.45)
    row_h = Inches(0.70)
    col1_w = Inches(2.10)
    col2_w = col_w - col1_w - Inches(0.10)
    for lab, body in left_rows:
        add_rect(s, x_l, row_y, col_w, row_h, fill=GREY_PALE, radius=0.04)
        add_textbox(s, x_l + Inches(0.15), row_y + Inches(0.08),
                    col1_w, row_h - Inches(0.16),
                    text=lab, size=11, bold=True, color=INK, line_spacing=1.20)
        add_textbox(s, x_l + Inches(0.15) + col1_w, row_y + Inches(0.08),
                    col2_w, row_h - Inches(0.16),
                    text=body, size=10, color=INK_2, line_spacing=1.25)
        row_y += row_h + Inches(0.05)

    # Правая — Низкий порог входа
    add_textbox(s, x_r, y0, col_w, Inches(0.36),
                text="НИЗКИЙ ПОРОГ ВХОДА (4 БАРЬЕРА СНЯТЫ)",
                size=10, bold=True, color=PRIMARY_DARK, font=FONT_DISPLAY)
    right_rows = [
        ("Технический", "1 команда /case-team; 0 шагов настройки multi-agent pipeline"),
        ("Доверия", "Дотренировка на 2 реальных кейс-чемпионатах; 100% цифр с URL"),
        ("Формата", "5 файлов на выходе: .xlsx, .pptx, SOLUTION.md, MANAGER-LOG.md, research"),
        ("Обучаемости", "0 команд изучения промпт-инжиниринга; задача — естественным языком"),
    ]
    row_y = y0 + Inches(0.45)
    row_h2 = Inches(0.88)
    for lab, body in right_rows:
        add_rect(s, x_r, row_y, col_w, row_h2, fill=PRIMARY_SOFT, radius=0.04)
        add_textbox(s, x_r + Inches(0.15), row_y + Inches(0.10),
                    col1_w, row_h2 - Inches(0.20),
                    text=lab, size=11, bold=True, color=PRIMARY_DARK,
                    line_spacing=1.20)
        add_textbox(s, x_r + Inches(0.15) + col1_w, row_y + Inches(0.10),
                    col2_w, row_h2 - Inches(0.20),
                    text=body, size=10, color=INK, line_spacing=1.30)
        row_y += row_h2 + Inches(0.05)

    draw_footer(s, 7)


def build_slide_08_team(prs):
    """Команда — карточка участника с placeholder-фото."""
    s = add_blank(prs)
    paint_background(s)
    draw_main_nav(s, SLIDE_TO_TAB[8])
    draw_header(s, 8,
                "Один продуктовый менеджер + команда из 5 ИИ-агентов = эффективная единица решения бизнес-кейсов",
                "Эта заявка собрана этой же командой /case-team за день до дедлайна — мета-доказательство тезиса.")

    # Левая часть — кружок-фото (placeholder) + label
    photo_d = Inches(3.0)
    photo_x = Inches(0.85)
    photo_y = Inches(3.30)

    # Проверка наличия файла фото
    photo_path = os.path.join(HERE, "assets", "team", "sergey.jpg")
    has_photo = os.path.exists(photo_path)

    if has_photo:
        # Вставить фото в круге (через cropped picture)
        s.shapes.add_picture(photo_path, photo_x, photo_y, photo_d, photo_d)
        # Обводка-кружок
        add_oval(s, photo_x, photo_y, photo_d, photo_d, line=PRIMARY)
    else:
        # Placeholder: круг primary_soft + инициалы СЛ
        add_oval(s, photo_x, photo_y, photo_d, photo_d, fill=PRIMARY_SOFT)
        add_textbox(s, photo_x, photo_y + Inches(0.55), photo_d, Inches(1.8),
                    text="СЛ", size=80, bold=True, color=PRIMARY_DARK,
                    font=FONT_DISPLAY, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(s, photo_x, photo_y + photo_d + Inches(0.10),
                    photo_d, Inches(0.30),
                    text="(фото в финальной версии)",
                    size=10, color=INK_2, align=PP_ALIGN.CENTER)

    # Правая часть — карточка с информацией
    x_info = Inches(4.40)
    y_info = Inches(3.10)
    info_w = Inches(8.40)

    add_textbox(s, x_info, y_info, info_w, Inches(0.65),
                text="Сергей Леонтьев", size=32, bold=True, color=INK,
                font=FONT_DISPLAY)
    add_textbox(s, x_info, y_info + Inches(0.70), info_w, Inches(0.40),
                text="Менеджер по продукту  ·  ООО «Нинтегра»",
                size=15, color=INK_2)
    add_textbox(s, x_info, y_info + Inches(1.15), info_w, Inches(0.40),
                text="МАИ 2025, «Прикладная информатика»  ·  Методология: AJTBD + ИИ-агенты",
                size=12, color=INK_2)

    # Тонкая линия-разделитель
    add_rect(s, x_info, y_info + Inches(1.70), Inches(0.60), Pt(2),
             fill=PRIMARY)

    # Достижения
    add_textbox(s, x_info, y_info + Inches(1.85), info_w, Inches(0.35),
                text="ДОСТИЖЕНИЯ", size=9, bold=True, color=PRIMARY_DARK,
                font=FONT_DISPLAY)
    add_textbox(s, x_info, y_info + Inches(2.18), info_w, Inches(0.85),
                text="Победитель, призёр и участник 10+ кейс-чемпионатов по бизнес-анализу, стратегическому консалтингу, продуктовому управлению и маркетингу.",
                size=13, color=INK, line_spacing=1.30)

    # Контакты — нижняя плашка во всю ширину
    cy = Inches(6.30)
    add_rect(s, Inches(0.55), cy, SLIDE_W - Inches(1.10), Inches(0.85),
             fill=PRIMARY_DARK, radius=0.06)
    add_textbox(s, Inches(0.85), cy + Inches(0.13),
                Inches(2.0), Inches(0.30),
                text="КОНТАКТЫ", size=10, bold=True, color=PRIMARY,
                font=FONT_DISPLAY, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, Inches(2.85), cy + Inches(0.10),
                Inches(4.0), Inches(0.32),
                text="suprrllysmanagement@mail.ru",
                size=12, color=WHITE, anchor=MSO_ANCHOR.TOP)
    add_textbox(s, Inches(2.85), cy + Inches(0.44),
                Inches(4.0), Inches(0.32),
                text="+7 (985) 176-06-58", size=12, color=WHITE)
    add_textbox(s, Inches(7.10), cy + Inches(0.10),
                Inches(5.7), Inches(0.32),
                text="GitHub: [github.com/Suprrllys/case-team](https://github.com/Suprrllys/case-team)",
                size=11, color=RGBColor(0xCC, 0xCC, 0xCC))
    add_textbox(s, Inches(7.10), cy + Inches(0.44),
                Inches(5.7), Inches(0.32),
                text="Видео-демо /case-team (5 мин): по запросу жюри",
                size=11, color=RGBColor(0xCC, 0xCC, 0xCC))

    draw_footer(s, 8)


def build_slide_09_final(prs):
    """Финал — hero accent:dark, формула differentiation."""
    s = add_blank(prs)
    paint_background(s, color=PRIMARY_DARK)

    add_rect(s, Inches(0.55), Inches(0.55), Inches(0.65), Inches(0.18),
             fill=PRIMARY)
    add_textbox(s, Inches(0.55), Inches(0.85), Inches(9.0), Inches(0.4),
                text="ФИНАЛ  ·  ТЕЗИС",
                size=11, bold=True, color=PRIMARY, font=FONT_DISPLAY)

    # Главный headline
    add_textbox(s, Inches(0.55), Inches(2.05), Inches(12.3), Inches(1.80),
                text="Полный пакет под защиту — в 8× быстрее, в стоимости подписки на нейросеть.",
                size=40, bold=True, color=WHITE, font=FONT_DISPLAY,
                line_spacing=1.15)

    # Подзаголовок-расшифровка
    add_textbox(s, Inches(0.55), Inches(4.30), Inches(12.3), Inches(1.30),
                text="Готовим полный пакет материалов решения бизнес-кейса под защиту в 8× быстрее и без дополнительной оплаты для пользователей Claude — за счёт мультиагентной AI-команды в Claude Code.",
                size=18, color=PRIMARY, line_spacing=1.40)

    # Slogan-line (цитата)
    add_rect(s, Inches(0.55), Inches(6.05), Inches(0.10), Inches(0.65),
             fill=PRIMARY)
    add_textbox(s, Inches(0.85), Inches(6.05), Inches(12.0), Inches(0.65),
                text="Альтернатива стратегическому консалтингу через мультиагентную AI-команду",
                size=15, italic=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    add_textbox(s, Inches(0.55), Inches(7.05), Inches(12.0), Inches(0.35),
                text="Сергей Леонтьев  ·  Product Manager  ·  Practicum Award 2026",
                size=11, color=RGBColor(0xCC, 0xCC, 0xCC))


# ────────────────────────────────────────────────────────────────────────
# APPENDICES A0..A9
# ────────────────────────────────────────────────────────────────────────

# Список 71 источника из YAML
SOURCES = [
    # Рынок консалтинга РФ (1-9)
    ("Рынок консалтинга РФ", [
        ("TAdviser «Консалтинг (рынок России)»", "https://www.tadviser.ru/index.php/Статья:Консалтинг_(рынок_России)", "общий объём 140 млрд ₽ в 2024 (+15% YoY)"),
        ("RAEX-RR рэнкинг стратегического консалтинга 2025", "https://raex-rr.com/b2b/consulting/consulting_strategic_planning_and_organizational_development/2025/analytics/", "стратегический сегмент 14.8 млрд ₽ (+39%)"),
        ("RAEX-RR крупнейшие консалт-группы 2025", "https://raex-rr.com/b2b/consulting/biggest_consulting_companies_and_groups/2025/", "Q1 2025 = 9.7 млрд ₽ (+9%)"),
        ("РБК Тренды «Тренды 2025 на рынке консалтинга»", "https://trends.rbc.ru/trends/social/67c0c6279a7947193f435524", "управленческий 37 млрд ₽ (+19%)"),
        ("HSE Daily «Рынок управленческого консалтинга в 2023»", "https://daily.hse.ru/post/novye-igroki-rynki-perspektivy-kak-vyglyadit-rynok-upravlencheskogo-konsaltinga-v-2023-godu", ""),
        ("Fless.pro «Большая тройка в России после 2022»", "https://fless.pro/big3-in-russia-after-2022-sbs-consulting", ""),
        ("TAdviser «Яков и Партнёры»", "https://www.tadviser.ru/index.php/Компания:Яков_и_Партнёры_(YnP,_ранее_McKinsey_в_России)", ""),
        ("KSK Group услуги стратегического консалтинга", "https://kskgroup.ru/services/consulting/strategicheskiy-konsalting/", "350-550 тыс ₽ — нижний сегмент"),
        ("Profi.ru стратегический консалтинг", "https://profi.ru/buhgaltery_i_yuristy/strategicheskii-konsalting/price/", ""),
    ]),
    # МСП РФ (10-13)
    ("МСП в РФ (TAM)", [
        ("Ведомости 25.12.2025 «Число субъектов МСП»", "https://www.vedomosti.ru/economics/articles/2025/12/25/1166484-chislo-subektov-msp-uvelichilos", "6.835 млн МСП в янв 2026"),
        ("Единый реестр МСП ФНС", "https://rmsp.nalog.ru/", "официальные данные"),
        ("Корпорация МСП «Численность обновила исторический максимум»", "https://xn--90aifddrld7a.xn--p1ai/news/country/chislennost-msp-v-rossii-obnovila-istoricheskiy-maksimum-i-prevysila-6-8-mln-predpriyatiy/", ""),
        ("Гарант «Количество малых и средних выросло»", "https://www.garant.ru/news/1964227/", ""),
    ]),
    # Ставки консультантов (14-18)
    ("Ставки консультантов и фрилансеров РФ", [
        ("dreamjob.ru «Старший консультант»", "https://dreamjob.ru/salary/starshiy-konsultant", "middle 150-280k ₽/мес"),
        ("Fless.pro «Зарплаты в Большой Тройке»", "https://fless.pro/consulting-salaries-ru", "Associate 350-450k, Manager 500k"),
        ("vc.ru «Стратегический консалтинг»", "https://vc.ru/hr/520245-strategicheskii-konsalting-na-blizhnem-vostoke-zarplaty-ot-120-tysyach-fantasticheskie-goroda-v-pustyne-i-metavselennaya", ""),
        ("Gorodrabot.ru «Зарплата консультанта в Москве»", "https://moskva.gorodrabot.ru/salaries/konsultant", ""),
        ("Lidopad «Фрилансеры в России 2025»", "https://lidopad.online/news/skolko-realno-zarabatyvayut-frilansery-v-rossii-statistika-i-prognoz-na-2026/", "средний доход 44k ₽/мес"),
    ]),
    # Multi-agent AI рынок (19-27)
    ("Multi-agent AI рынок и тренды", [
        ("Gartner «40% enterprise apps with AI agents by 2026»", "https://www.gartner.com/en/newsroom/press-releases/2025-08-26-gartner-predicts-40-percent-of-enterprise-apps-will-feature-task-specific-ai-agents-by-2026-up-from-less-than-5-percent-in-2025", ""),
        ("McKinsey MGI «$4.4 trillion AI value»", "https://www.mckinsey.com/mgi/media-center/ai-could-increase-corporate-profits-by-4-trillion-a-year-according-to-new-research", ""),
        ("McKinsey «State of AI» ноябрь 2025", "https://www.mckinsey.com/capabilities/quantumblack/our-insights/the-state-of-ai", ""),
        ("McKinsey «Economic potential of generative AI»", "https://www.mckinsey.com/capabilities/tech-and-ai/our-insights/the-economic-potential-of-generative-ai-the-next-productivity-frontier", ""),
        ("Markets&Markets «AI Agents Market 2025-2030»", "https://www.marketsandmarkets.com/Market-Reports/ai-agents-market-15761548.html", ""),
        ("Paul Okhrem «Enterprise AI Agents Statistics 2026»", "https://paul-okhrem.com/enterprise-ai-agents-statistics-2026/", "$7.6→$10.8 млрд + 40% projects cancelled"),
        ("Joget «AI Agent Adoption 2026»", "https://joget.com/ai-agent-adoption-in-2026-what-the-analysts-data-shows/", ""),
        ("a16z «How 100 Enterprise CIOs Build Gen AI 2025»", "https://a16z.com/ai-enterprise-2025/", ""),
        ("a16z «Where Enterprises Adopt AI»", "https://a16z.com/where-enterprises-are-actually-adopting-ai/", ""),
    ]),
    # Конкуренты int (28-41)
    ("Конкуренты — международные multi-agent", [
        ("Felloai «Manus AI Pricing 2026»", "https://felloai.com/manus-ai-pricing/", "$0/$19/$39/$199"),
        ("RioTimes «Manus AI 14 Failures»", "https://www.riotimesonline.com/manus-a-i-review-14-failures-in-two-weeks-of-testing/", ""),
        ("Trickle «Manus AI Review»", "https://trickle.so/blog/manus-ai-review", "Reddit credit-burn $2380"),
        ("Trustpilot Manus AI reviews", "https://www.trustpilot.com/review/manus-ai.sbs", ""),
        ("Lindy.ai Pricing", "https://www.lindy.ai/pricing", "Free / $19.99 / $49.99"),
        ("CloudTalk «Lindy AI Pricing 2026»", "https://www.cloudtalk.io/blog/lindy-ai-pricing/", "onboarding $1500 hidden"),
        ("CrewAI Pricing", "https://crewai.com/pricing", ""),
        ("ZenML «CrewAI Pricing Guide»", "https://www.zenml.io/blog/crewai-pricing", "Enterprise $60-120k/год"),
        ("Microsoft Research «Magentic-One»", "https://www.microsoft.com/en-us/research/articles/magentic-one-a-generalist-multi-agent-system-for-solving-complex-tasks/", ""),
        ("WinBuzzer «AutoGen 0.4 + Magentic-One»", "https://winbuzzer.com/2025/01/14/microsoft-releases-autogen-0-4-with-magentic-one-multi-ai-agent-framework-xcxwbn/", ""),
        ("Devin Pricing", "https://devin.ai/pricing/", "Core $20 / Team $500"),
        ("VentureBeat «Devin 2.0 Slashes Price»", "https://venturebeat.com/programming-development/devin-2-0-is-here-cognition-slashes-price-of-ai-software-engineer-to-20-per-month-from-500", ""),
        ("OpenAIToolsHub «Devin AI Review: 13.86% SWE-Bench»", "https://www.openaitoolshub.org/en/blog/devin-ai-review", ""),
        ("SitePoint «Devin Aftermath»", "https://www.sitepoint.com/devin-ai-engineers-production-realities/", ""),
    ]),
    # Российские AI (42-49)
    ("Конкуренты — российские AI-платформы", [
        ("GigaChat Enterprise на Хабре", "https://habr.com/ru/companies/sberbank/news/1005990/", "запуск март 2026"),
        ("vc.ru «Сбер GigaChat Enterprise»", "https://vc.ru/typespace/2769088-sber-gigachat-enterprise-platforma-dlya-ii-agentov", ""),
        ("GigaChat Business — Multi-Agent System", "https://b2b.giga.chat/multi-agent-system", ""),
        ("Yandex AI Studio главная", "https://aistudio.yandex.ru/ru", ""),
        ("Яндекс «AI Studio большое обновление»", "https://yandex.ru/company/news/03-03-2026-01", ""),
        ("Habr «Яндекс открыл рассуждающие ИИ-агенты»", "https://habr.com/ru/amp/publications/1009872/", ""),
        ("Just AI Agent Platform", "https://agentplatform.just-ai.com/", ""),
        ("Just AI «Сравнение российских AI-платформ 2026»", "https://just-ai.com/blog/sravnenie-rossijskih-platform-dlya-sozdaniya-ai-agentov", ""),
    ]),
    # Универсальные AI-чаты (50-53)
    ("Универсальные AI-чаты + презентации", [
        ("Claude for PowerPoint", "https://claude.com/claude-for-powerpoint", ""),
        ("Prezent.ai «Claude for PowerPoint»", "https://www.prezent.ai/blog/claude-for-powerpoint", ""),
        ("MindStudio «Gamma vs ChatGPT vs Claude vs Google»", "https://www.mindstudio.ai/blog/gamma-vs-chatgpt-vs-claude-vs-google-slides-ai-presentation-tool-comparison", ""),
        ("SlideSpeak «Create Presentations with Claude Design 2026»", "https://slidespeak.co/blog/create-presentations-claude-design", ""),
    ]),
    # Бенчмарки (54-62)
    ("Бенчмарки AI-экономии", [
        ("Klarna Press «AI assistant 2/3 customer service»", "https://www.klarna.com/international/press/klarna-ai-assistant-handles-two-thirds-of-customer-service-chats-in-its-first-month/", ""),
        ("OpenAI «Klarna AI does work of 700 agents»", "https://openai.com/index/klarna/", ""),
        ("CX Dive «Klarna AI slashing costs»", "https://www.customerexperiencedive.com/news/klarna-ai-slash-customer-service-costs/748647/", ""),
        ("Twig «Klarna AI Cut Resolution Time 82%»", "https://www.twig.so/blog/how-klarna-is-revolutionizing-customer-support-with-ai", ""),
        ("GitHub Blog «Quantifying GitHub Copilot Impact»", "https://github.blog/news-insights/research/research-quantifying-github-copilots-impact-on-developer-productivity-and-happiness/", "+55% speed"),
        ("arXiv «AI on Developer Productivity: Copilot»", "https://arxiv.org/abs/2302.06590", "P=0.0017"),
        ("Anthropic «How AI Is Transforming Work at Anthropic»", "https://www.anthropic.com/research/how-ai-is-transforming-work-at-anthropic", ""),
        ("Anthropic «Estimating productivity gains»", "https://www.anthropic.com/research/estimating-productivity-gains", "~80-81% task speedup"),
        ("HUB International «Anthropic Claude to 20 000+»", "https://www.hubinternational.com/media-center/press-releases/2026/02/hub-international-brings-anthropics-claude-to-20000-employees/", ""),
    ]),
    # Финмодель (63-66)
    ("Финмодель — стоимость API + курс ЦБ РФ", [
        ("Anthropic Pricing", "https://platform.claude.com/docs/en/about-claude/pricing", "Opus 4.7 $5/$25 на 1M токенов"),
        ("Finout «Anthropic API Pricing 2026»", "https://www.finout.io/blog/anthropic-api-pricing", ""),
        ("ЦБ РФ официальный курс", "https://www.cbr.ru/", "USD/RUB 73.34 на 14 мая 2026"),
        ("Finance.mail.ru «Курс доллара на 14 мая 2026»", "https://finance.mail.ru/article/oficialnyj-kurs-dollara-na-14-maya-2026-69208634/", ""),
    ]),
    # Методологические источники (67-71)
    ("Методологические источники", [
        ("Иван Замесин — публикации по AJTBD-консалтингу", "https://zamesin.ru/", "Mot-Moment, extract-jobs-from-reviews, и др."),
        ("Christensen «Competing Against Luck» (JTBD framework)", "https://www.amazon.com/Competing-Against-Luck-Innovation-Customer/dp/0062435612", ""),
        ("Osterwalder «Value Proposition Design»", "https://www.strategyzer.com/library/value-proposition-design", ""),
        ("Anthropic Claude Design — Claude.ai", "https://claude.com", ""),
        ("Practicum Award 2026 — лендинг конкурса", "https://practicum.yandex.ru/practicum-award", "критерии оценки + жюри"),
    ]),
]


def build_appendix_header(s, slide_num, app_index, total_apps, title, subtitle):
    paint_background(s)
    draw_appendix_nav(s, app_index, total_apps)
    add_textbox(s, Inches(0.55), Inches(0.85), SLIDE_W - Inches(1.10),
                Inches(1.05), text=title, size=22, bold=True, color=INK,
                font=FONT_DISPLAY, line_spacing=1.10)
    if subtitle:
        add_textbox(s, Inches(0.55), Inches(2.00), SLIDE_W - Inches(1.10),
                    Inches(0.55), text=subtitle, size=12, color=INK_2,
                    line_spacing=1.30)
    add_rect(s, Inches(0.55), Inches(2.65), Inches(0.60), Pt(2),
             fill=PRIMARY)


def _add_sources_block(slide, x, y, w, h, items, start_num):
    """Рендер блока источников в textbox с гиперссылками."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    # Первая строка инициализирует параграф
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    p0.line_spacing = 1.20
    first = True
    n = start_num
    for title, url, ctx in items:
        if first:
            p = p0
            first = False
        else:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.20
            p.space_before = Pt(2)
        # Номер
        add_run(p, f"{n}. ", size=9, color=INK_2, font=FONT_BODY)
        # Кликабельный заголовок
        add_run(p, title, size=9, color=LINK_BLUE, font=FONT_BODY,
                hyperlink=url)
        if ctx:
            add_run(p, f"  — {ctx}", size=9, color=INK_2, font=FONT_BODY,
                    italic=True)
        n += 1
    return n


def build_appendix_a0_sources(prs):
    """A0 — Источники (71 шт). Разбит на 3 слайда: A0.1, A0.2, A0.3."""
    flat = []  # все 71 источника как плоский список с категорией
    for cat, items in SOURCES:
        for title, url, ctx in items:
            flat.append((cat, title, url, ctx))

    # 71 источник: разбьём на 3 части ~24 на слайд (по 8-10 шт в колонке × 3 колонки)
    chunk_size = 24
    chunks = [flat[i:i + chunk_size] for i in range(0, len(flat), chunk_size)]
    # Получится ~3 чанка

    for ci, chunk in enumerate(chunks):
        s = add_blank(prs)
        slide_num = 10 + ci  # A0 — это слайд 10, 11, 12 в общей нумерации (если 3 части)
        title = f"Источники ({ci * chunk_size + 1}–{ci * chunk_size + len(chunk)} из 71)"
        sub = "Все источники, использованные при формировании решения — публичные URL, каждый кликабельный."
        build_appendix_header(s, slide_num, 0, 9, title, sub)

        # 3 колонки источников
        col_w = (SLIDE_W - Inches(1.10) - Inches(0.40)) / 3  # 3 колонки с gap
        col_gap = Inches(0.20)
        col_h = Inches(4.30)
        col_y = Inches(2.85)
        # Распределить chunk на 3 колонки приблизительно равномерно
        per_col = (len(chunk) + 2) // 3
        cols = [chunk[i*per_col:(i+1)*per_col] for i in range(3)]
        # Хвост, если есть
        used = sum(len(c) for c in cols)
        if used < len(chunk):
            cols[-1].extend(chunk[used:])

        cur_n = ci * chunk_size + 1
        cx = Inches(0.55)
        for col_items in cols:
            # Передадим только (title, url, ctx) — игнорируем категорию здесь
            items = [(t, u, c) for (_, t, u, c) in col_items]
            new_n = _add_sources_block(s, cx, col_y, col_w, col_h, items, cur_n)
            cur_n = new_n
            cx += col_w + col_gap

        draw_footer(s, slide_num)


def build_appendix_simple(prs, slide_num, app_index, title, sections,
                          subtitle=""):
    """Универсальный шаблон приложения: список заголовок→буллеты.
    Если секций >3 — авто-разбивка на 2 колонки.
    """
    s = add_blank(prs)
    build_appendix_header(s, slide_num, app_index, 9, title, subtitle)

    use_2col = len(sections) > 3

    if use_2col:
        # Раскидать секции на 2 колонки приблизительно равномерно
        half = (len(sections) + 1) // 2
        col1_secs = sections[:half]
        col2_secs = sections[half:]
        col_w = (SLIDE_W - Inches(1.10) - Inches(0.30)) / 2
        x_left = Inches(0.55)
        x_right = x_left + col_w + Inches(0.30)
        for col_secs, x_col in [(col1_secs, x_left), (col2_secs, x_right)]:
            y = Inches(2.95)
            for sec_title, bullets in col_secs:
                add_textbox(s, x_col, y, col_w, Inches(0.30),
                            text=sec_title.upper(), size=9, bold=True,
                            color=PRIMARY_DARK, font=FONT_DISPLAY)
                y += Inches(0.30)
                # Каждый буллет ~0.36" при size 10
                bh = Inches(0.34) * len(bullets) + Inches(0.05)
                tb = add_textbox(s, x_col + Inches(0.15), y,
                                 col_w - Inches(0.15), bh,
                                 text="", size=10, color=INK,
                                 line_spacing=1.30)
                tf = tb.text_frame
                first = True
                for b in bullets:
                    if first:
                        set_text(tf, "•  " + b, size=10, color=INK,
                                 line_spacing=1.30)
                        first = False
                    else:
                        add_paragraph(tf, "•  " + b, size=10, color=INK,
                                      line_spacing=1.30, space_before=2)
                y += bh + Inches(0.12)
    else:
        y = Inches(2.95)
        x = Inches(0.55)
        w = SLIDE_W - Inches(1.10)
        for sec_title, bullets in sections:
            add_textbox(s, x, y, w, Inches(0.34),
                        text=sec_title.upper(), size=10, bold=True,
                        color=PRIMARY_DARK, font=FONT_DISPLAY)
            y += Inches(0.36)
            bh = Inches(0.34) * len(bullets) + Inches(0.10)
            tb = add_textbox(s, x + Inches(0.20), y, w - Inches(0.20), bh,
                             text="", size=11, color=INK, line_spacing=1.30)
            tf = tb.text_frame
            first = True
            for b in bullets:
                if first:
                    set_text(tf, "•  " + b, size=11, color=INK,
                             line_spacing=1.30)
                    first = False
                else:
                    add_paragraph(tf, "•  " + b, size=11, color=INK,
                                  line_spacing=1.30, space_before=3)
            y += bh + Inches(0.16)

    draw_footer(s, slide_num)


def build_appendix_a1_segment(prs, slide_num):
    s = add_blank(prs)
    build_appendix_header(s, slide_num, 1, 9,
        "Сегмент клиента — фаундер МСП в стратегическом моменте",
        "Founder/CEO Series A-B или Operations / Commercial Director МСП в РФ; штат 5-50; оборот 10-500 млн ₽/год; города-миллионники.")

    # 3 колонки: иерархия работ / триггеры / критерии успеха
    y0 = Inches(2.95)
    col_w = Inches(4.05)
    col_gap = Inches(0.10)
    col_h = Inches(3.85)

    # Колонка 1 — Иерархия работ
    x1 = Inches(0.55)
    add_rect(s, x1, y0, col_w, col_h, fill=PRIMARY_SOFT, radius=0.06)
    add_textbox(s, x1 + Inches(0.20), y0 + Inches(0.15), col_w - Inches(0.40),
                Inches(0.32), text="ИЕРАРХИЯ РАБОТ (AJTBD)",
                size=10, bold=True, color=PRIMARY_DARK, font=FONT_DISPLAY)
    tb = add_textbox(s, x1 + Inches(0.20), y0 + Inches(0.50),
                     col_w - Inches(0.40), col_h - Inches(0.70),
                     text="", size=10, color=INK, line_spacing=1.30)
    tf = tb.text_frame
    set_text(tf, "Big Job (верхний уровень):", size=10, bold=True,
             color=PRIMARY_DARK, line_spacing=1.25)
    add_paragraph(tf, "Принять обоснованное стратегическое решение и провести его через стейкхолдеров.",
                  size=10, color=INK, line_spacing=1.30, space_before=2)
    add_paragraph(tf, "Core Job (выполняется /case-team):",
                  size=10, bold=True, color=PRIMARY_DARK,
                  line_spacing=1.25, space_before=8)
    add_paragraph(tf, "Подготовить полный набор материалов под защиту.",
                  size=10, color=INK, line_spacing=1.30, space_before=2)
    add_paragraph(tf, "Small Jobs (остаются за человеком):",
                  size=10, bold=True, color=PRIMARY_DARK,
                  line_spacing=1.25, space_before=8)
    for it in ["go/no-go решение", "защита лично перед стейкхолдерами",
               "внедрить стратегию операционно", "договориться с партнёрами"]:
        add_paragraph(tf, "•  " + it, size=10, color=INK, line_spacing=1.25,
                      space_before=2)

    # Колонка 2 — Триггеры
    x2 = x1 + col_w + col_gap
    add_rect(s, x2, y0, col_w, col_h, fill=GREY_PALE, radius=0.06)
    add_textbox(s, x2 + Inches(0.20), y0 + Inches(0.15), col_w - Inches(0.40),
                Inches(0.32), text="ТРИГГЕРЫ CORE JOB",
                size=10, bold=True, color=PRIMARY_DARK, font=FONT_DISPLAY)
    triggers = [
        ("Питч инвестору в среду", "7-14 дней"),
        ("Совет директоров в пятницу", "3-10 дней"),
        ("Конкурент сделал X", "1-3 дня"),
        ("Go/no-go нового продукта", "7-21 день"),
    ]
    ty = y0 + Inches(0.60)
    for label, deadline in triggers:
        add_textbox(s, x2 + Inches(0.20), ty, col_w - Inches(0.40),
                    Inches(0.30), text=label, size=11, bold=True, color=INK)
        add_textbox(s, x2 + Inches(0.20), ty + Inches(0.30),
                    col_w - Inches(0.40), Inches(0.28),
                    text="дедлайн: " + deadline, size=10, color=INK_2,
                    italic=False)
        ty += Inches(0.70)
    add_textbox(s, x2 + Inches(0.20), ty, col_w - Inches(0.40),
                Inches(0.60),
                text="Частотность работы: 2-6 раз в год — [граф работ из отзывов клиентов на конкурентов](https://github.com/Suprrllys/case-team/blob/main/solution/analyst/jobs-and-segments-analysis.md).",
                size=10, italic=True, color=INK_2, line_spacing=1.25)

    # Колонка 3 — Критерии успеха
    x3 = x2 + col_w + col_gap
    add_rect(s, x3, y0, col_w, col_h, fill=PRIMARY_SOFT, radius=0.06)
    add_textbox(s, x3 + Inches(0.20), y0 + Inches(0.15), col_w - Inches(0.40),
                Inches(0.32), text="КРИТЕРИИ УСПЕХА CORE JOB",
                size=10, bold=True, color=PRIMARY_DARK, font=FONT_DISPLAY)
    crits = [
        ("Скорость", "4 часа vs 5-6 дней"),
        ("Прослеживаемость", "100% цифр с URL-источником"),
        ("Защищаемость Q&A", "1 цифра = 1 источник"),
        ("Воспроизводимость", "5 файлов на выходе"),
        ("Порог входа", "1 команда, 0 настройки"),
    ]
    cy = y0 + Inches(0.55)
    for label, val in crits:
        add_textbox(s, x3 + Inches(0.20), cy, col_w - Inches(0.40),
                    Inches(0.26), text=label, size=10, bold=True, color=INK)
        add_textbox(s, x3 + Inches(0.20), cy + Inches(0.26),
                    col_w - Inches(0.40), Inches(0.30),
                    text=val, size=11, color=ACCENT_POS, bold=True,
                    font=FONT_DISPLAY)
        cy += Inches(0.62)

    # Нижняя сноска — главная эмоция клиента
    note_y = Inches(6.90)
    add_textbox(s, Inches(0.55), note_y, SLIDE_W - Inches(1.10), Inches(0.40),
                text="ГЛАВНЫЙ СТРАХ КЛИЕНТА — не «дорого», а «обнаружится фейк под вопросами стейкхолдера». /case-team функционально замещает «защиту через репутацию» через прозрачность методологии.",
                size=10, color=INK_2, italic=True, line_spacing=1.25)

    draw_footer(s, slide_num)


def build_appendix_a2_competitors(prs, slide_num):
    s = add_blank(prs)
    build_appendix_header(s, slide_num, 2, 9,
        "Конкуренты в multi-agent — матрица позиционирования",
        "9 проанализированных конкурентов по 3 слоям. Мы занимаем пустую клетку в матрице — готовый продукт под стратегический класс задач для одного человека без IT-ресурса.")

    # 3 группы конкурентов в колонках
    y0 = Inches(2.95)
    col_w = Inches(4.05)
    col_gap = Inches(0.10)
    col_h = Inches(3.85)

    groups = [
        ("ПРЯМЫЕ (multi-agent для бизнес-задач)", PRIMARY_SOFT, [
            ("Manus AI", "autonomous generalist; мы = команда ролей с куратором смысла"),
            ("Lindy", "операционная рутина; мы = стратегические задачи"),
            ("CrewAI / AutoGen", "конструкторы для разработчиков; мы = готовый продукт"),
        ]),
        ("РОССИЙСКИЕ (operational, не strategic)", GREY_PALE, [
            ("GigaChat Enterprise (Сбер)", "для крупных корпов"),
            ("Yandex AI Studio", "конструктор для разработчиков"),
            ("Just AI Agent Platform", "Jay Knowledge Hub + Jay Flow"),
        ]),
        ("КОСВЕННЫЕ", PRIMARY_SOFT, [
            ("Cognition Devin", "single-role SWE; 85% failure rate без супервизии"),
            ("ChatGPT / Claude.ai", "главный когнитивный конкурент"),
            ("Big4 + «Яков и партнёры»", "нанимаются за защиту от позора, не за работу"),
            ("Фрилансеры с биржи", "никто не сшивает в питч-док"),
        ]),
    ]

    for i, (gtitle, gfill, items) in enumerate(groups):
        x = Inches(0.55) + i * (col_w + col_gap)
        add_rect(s, x, y0, col_w, col_h, fill=gfill, radius=0.06)
        add_textbox(s, x + Inches(0.20), y0 + Inches(0.15),
                    col_w - Inches(0.40), Inches(0.32),
                    text=gtitle, size=10, bold=True, color=PRIMARY_DARK,
                    font=FONT_DISPLAY)
        cy = y0 + Inches(0.55)
        for name, body in items:
            add_textbox(s, x + Inches(0.20), cy, col_w - Inches(0.40),
                        Inches(0.30), text=name, size=12, bold=True, color=INK,
                        font=FONT_DISPLAY)
            add_textbox(s, x + Inches(0.20), cy + Inches(0.30),
                        col_w - Inches(0.40), Inches(0.55),
                        text=body, size=10, color=INK_2, line_spacing=1.25)
            cy += Inches(0.85)

    note_y = Inches(6.90)
    add_textbox(s, Inches(0.55), note_y, SLIDE_W - Inches(1.10), Inches(0.40),
                text="Полная матрица 8 критериев × 9 конкурентов — в [открытом репозитории](https://github.com/Suprrllys/case-team/blob/main/solution/analyst/practicum-award-context.md).",
                size=10, italic=True, color=INK_2, line_spacing=1.25)
    draw_footer(s, slide_num)


def build_appendix_a3_tech(prs, slide_num):
    sections = [
        ("Модель", ["Claude Opus 4.7 (1M context) — длинный контекст для многошаговых кейсов; state-of-art reasoning."]),
        ("Архитектура", ["Claude Code subagents — изолированные роли с собственными промптами, инструментами, папками."]),
        ("Интеграции", ["MCP-серверы (Notion, Google Drive)", "Web search / fetch — доступ к внешним источникам без копирования вручную."]),
        ("База знаний", ["Локальные эталоны (Cup Moscow + Inno)", "Универсальные инструкции (финмодель, презентация)", "Каталог 18 типов слайдов + каталог исследовательских фреймворков."]),
        ("Методологии", ["AJTBD-промпты Замесина (market-trends, competitor-research, differentiation-strategy, extract-jobs-from-reviews) интегрированы в case-marketer и case-analyst."]),
        ("Дизайн-качество", ["Скилл claude-design + 14 фазовых процедур (ai-slop, accessibility, hierarchy-rhythm, polish)."]),
        ("Сборка .pptx", ["python-pptx + pre-flight overflow check + macOS PowerPoint blessing — воспроизводимая программная сборка."]),
    ]
    build_appendix_simple(prs, slide_num, 3,
                          "Технологический стек",
                          sections,
                          subtitle="Каждый компонент стека — обоснованный выбор под класс задачи «полный пакет под защиту за 4 часа».")


def build_appendix_a4_finmodel(prs, slide_num):
    s = add_blank(prs)
    build_appendix_header(s, slide_num, 4, 9,
        "Детальная финмодель и sensitivity",
        "Реальные эмпирические данные из решений для Axenix и Askona; ставки и стоимость API — на дату 14 мая 2026.")

    # 3 блока: baseline / ставки / стоимость
    y0 = Inches(2.95)
    col_w = Inches(4.05)
    col_gap = Inches(0.10)
    col_h = Inches(3.60)

    # Блок 1 — Baseline
    x1 = Inches(0.55)
    add_rect(s, x1, y0, col_w, col_h, fill=PRIMARY_SOFT, radius=0.06)
    add_textbox(s, x1 + Inches(0.20), y0 + Inches(0.15), col_w - Inches(0.40),
                Inches(0.32), text="ЭМПИРИЧЕСКИЕ ДАННЫЕ (BASELINE)",
                size=10, bold=True, color=PRIMARY_DARK, font=FONT_DISPLAY)
    items = [
        ("Cup Moscow 2024 (Axenix)", "5 × 140 чел.-часов"),
        ("Inno Case Hack 2025 (Askona)", "5 × 40 чел.-часов"),
        ("Средний кейс", "~90 чел.-часов команды"),
    ]
    cy = y0 + Inches(0.55)
    for lab, val in items:
        add_textbox(s, x1 + Inches(0.20), cy, col_w - Inches(0.40),
                    Inches(0.30), text=lab, size=11, bold=True, color=INK)
        add_textbox(s, x1 + Inches(0.20), cy + Inches(0.30),
                    col_w - Inches(0.40), Inches(0.30),
                    text=val, size=12, color=PRIMARY_DARK, bold=True,
                    font=FONT_DISPLAY)
        cy += Inches(0.75)

    # Блок 2 — Ставки
    x2 = x1 + col_w + col_gap
    add_rect(s, x2, y0, col_w, col_h, fill=GREY_PALE, radius=0.06)
    add_textbox(s, x2 + Inches(0.20), y0 + Inches(0.15), col_w - Inches(0.40),
                Inches(0.32), text="СТАВКИ КОМАНДЫ (MIDDLE, ₽/Ч)",
                size=10, bold=True, color=PRIMARY_DARK, font=FONT_DISPLAY)
    rates = [
        ("PM", "1 190"),
        ("Аналитик", "1 070"),
        ("Финансист", "1 190"),
        ("Маркетолог", "1 010"),
        ("Дизайнер", "950"),
        ("Средневзвешенная", "1 100"),
    ]
    cy = y0 + Inches(0.55)
    for lab, val in rates:
        bold = (lab == "Средневзвешенная")
        col = ACCENT_POS if bold else INK
        add_textbox(s, x2 + Inches(0.20), cy, col_w - Inches(1.10),
                    Inches(0.28), text=lab, size=11, bold=bold, color=INK)
        add_textbox(s, x2 + col_w - Inches(1.20), cy, Inches(1.0),
                    Inches(0.28), text=val, size=11, bold=True, color=col,
                    font=FONT_DISPLAY, align=PP_ALIGN.RIGHT)
        cy += Inches(0.35)

    # Блок 3 — Стоимость 1 кейса
    x3 = x2 + col_w + col_gap
    add_rect(s, x3, y0, col_w, col_h, fill=PRIMARY_DARK, radius=0.06)
    add_textbox(s, x3 + Inches(0.20), y0 + Inches(0.15), col_w - Inches(0.40),
                Inches(0.32), text="СТОИМОСТЬ 1 КЕЙСА",
                size=10, bold=True, color=PRIMARY, font=FONT_DISPLAY)
    cost = [
        ("С агентами (внутр.)", "10 930 ₽", ACCENT_POS),
        ("In-house из 5", "99 000 ₽", WHITE),
        ("Big4-tier (×4)", "396 000 ₽", RGBColor(0xFF, 0xC8, 0xC8)),
    ]
    cy = y0 + Inches(0.55)
    for lab, val, col in cost:
        add_textbox(s, x3 + Inches(0.20), cy, col_w - Inches(0.40),
                    Inches(0.30), text=lab, size=11, color=RGBColor(0xBB, 0xBB, 0xBB))
        add_textbox(s, x3 + Inches(0.20), cy + Inches(0.30),
                    col_w - Inches(0.40), Inches(0.42),
                    text=val, size=22, bold=True, color=col, font=FONT_DISPLAY)
        cy += Inches(0.85)
    add_textbox(s, x3 + Inches(0.20), cy, col_w - Inches(0.40),
                Inches(0.40),
                text="Для клиента (D10 подача А): $0 — Claude подписка уже оплачена.",
                size=10, color=PRIMARY, italic=True, line_spacing=1.25)

    # Footer-note: sensitivity
    note_y = Inches(6.75)
    add_rect(s, Inches(0.55), note_y, SLIDE_W - Inches(1.10), Inches(0.50),
             fill=PRIMARY_SOFT, radius=0.06)
    add_textbox(s, Inches(0.75), note_y + Inches(0.10), SLIDE_W - Inches(1.50),
                Inches(0.30),
                text="SENSITIVITY ±20% — устойчиво 84-98% vs in-house. Источники — A0. [Полная финмодель](https://github.com/Suprrllys/case-team/blob/main/solution/financier/unit-economics.md) в открытом репозитории.",
                size=11, color=PRIMARY_DARK, line_spacing=1.25,
                anchor=MSO_ANCHOR.MIDDLE)

    draw_footer(s, slide_num)


def build_appendix_a5_cascade(prs, slide_num):
    sections = [
        ("Уровень 1 — Бизнес-цель", [
            "Выход в финалисты Practicum Award + борьба за приз номинации (500k ₽ + грант Yandex AI Studio) или спецприз «Будущее за ним» (150k ₽).",
        ]),
        ("Уровень 2 — Маркетинговые цели (восприятие жюри)", [
            "M1: запоминание как «новый класс инструмента», не «ещё одна multi-agent»",
            "M2: вера в цифры экономии (84-98%) без «слишком красиво»",
            "M3: восприятие как методологически зрелой заявки",
            "M4: ассоциация с продуктовой и методологической культурой РФ",
            "M5: мета-демо — уникальный аргумент, не маркетинговый трюк",
        ]),
        ("Уровень 3 — Коммуникационные цели", [
            "C1 (Slide 2): «всё решение за 30 секунд через 3 цифры»",
            "C2 (Slide 3): «5 текущих способов — дороги, медленны или без артефактов»",
            "C3 (Slide 4): «5 ролей с куратором смысла, физические артефакты, эталоны»",
            "C4 (Slide 6): «категориальный сдвиг — в стоимости подписки на нейросеть»",
            "C5 (Slide 7): «главный страх клиента — не дорого, а фейк под вопросами»",
        ]),
        ("Уровень 4 — Метрики успеха", [
            "ME1: все 8 членов жюри Техно-номинации читают заявку целиком",
            "ME2: one-liner «в стоимости подписки» повторяется на 5 слайдах",
            "ME3: попадание в топ-10 финалистов",
            "ME4: brand-lift — упоминание категории в обратной связи жюри",
        ]),
    ]
    build_appendix_simple(prs, slide_num, 5,
                          "Каскад целей маркетинга (4 уровня)", sections,
                          subtitle="От бизнес-цели до измеримых метрик восприятия жюри — адаптация Reach / Frequency / SOV под формат конкурса.")


def build_appendix_a6_roadmap(prs, slide_num):
    s = add_blank(prs)
    build_appendix_header(s, slide_num, 6, 9,
        "Дорожная карта — Now / Next / Later",
        "Текущая зрелость, ближайшие итерации и долгосрочная траектория продукта /case-team.")

    # 3 колонки Now / Next / Later
    y0 = Inches(2.95)
    col_w = Inches(4.05)
    col_gap = Inches(0.10)
    col_h = Inches(4.05)

    cols = [
        ("NOW · готово", PRIMARY, [
            "5 агентов с явными ролями",
            "2 эталона дотренировки (Cup Moscow + Inno) → 26 логик",
            "2 режима handoff (premium HTML / regular .pptx)",
            "Скилл /case-team — одна команда установки",
            "Реальная подача в Practicum Award собрана этой командой",
        ]),
        ("NEXT · 1-3 мес", PRIMARY_SOFT, [
            "Публикация агентов и skill целиком в [github.com/Suprrllys/case-team](https://github.com/Suprrllys/case-team) (сейчас опубликованы материалы заявки)",
            "Расширение эталонной базы (антикризис, M&A, креативный бренд)",
            "3-5 pilot-запусков с B2B-клиентами в обмен на кейс-стади",
        ]),
        ("LATER · 3-12 мес", GREY_PALE, [
            "Полноценный B2B-go-to-market (стартапы + малый бизнес + L&D)",
            "Edu-канал через бизнес-школы и кейс-чемпионаты",
            "Возможная PRO-надстройка на базе Anthropic API",
        ]),
    ]

    for i, (title, fill, items) in enumerate(cols):
        x = Inches(0.55) + i * (col_w + col_gap)
        add_rect(s, x, y0, col_w, col_h, fill=fill, radius=0.08)
        # Заголовок
        title_color = PRIMARY_DARK if i == 0 else PRIMARY_DARK
        add_textbox(s, x + Inches(0.25), y0 + Inches(0.20), col_w - Inches(0.50),
                    Inches(0.40), text=title, size=14, bold=True,
                    color=title_color, font=FONT_DISPLAY)
        # Bullets
        bullets_h = col_h - Inches(0.75)
        tb = add_textbox(s, x + Inches(0.25), y0 + Inches(0.65),
                         col_w - Inches(0.50), bullets_h,
                         text="", size=11, color=INK, line_spacing=1.35)
        tf = tb.text_frame
        first = True
        for it in items:
            if first:
                set_text(tf, "•  " + it, size=11, color=INK, line_spacing=1.35)
                first = False
            else:
                add_paragraph(tf, "•  " + it, size=11, color=INK,
                              line_spacing=1.35, space_before=5)

    draw_footer(s, slide_num)


def build_appendix_a7_methods(prs, slide_num):
    sections = [
        ("1. Эталоны команды", [
            "Cup Moscow 2024 (Axenix, авиахаб Звартноц) — стратегический B2G",
            "Inno Case Hack 2025 (Askona, A-Sleep) — стартап-маркетинговый",
            "Извлечено 26 методологических логик из этих решений",
        ]),
        ("2a. Бизнес-классика — JTBD-блок", [
            "AJTBD / JTBD (Christensen, Moesta) — сегментация по работам",
            "Mot-Moment Persona — JTBD-консалтинг (Замесин)",
            "AJTBD-промпты Замесина в case-marketer и case-analyst",
            "Value Proposition Canvas (Osterwalder)",
        ]),
        ("2b. Бизнес-классика — приоритизация и процессы", [
            "MoSCoW / RICE / Weighted Scoring — приоритизация",
            "AS-IS / TO-BE — Lean / Six Sigma process improvement",
            "Theory of Change — impact-инвестирование",
            "Big4-tier consulting structure (presale + методология)",
        ]),
        ("3. Эвристики автора", [
            "Бюджет слайда (1 тезис + 3-7 буллетов + 1 визуал)",
            "Правило 50:50 текст/визуал на контентных слайдах",
            "Pre-flight overflow check для python-pptx",
        ]),
    ]
    build_appendix_simple(prs, slide_num, 7,
                          "Методологические источники (иерархия)", sections,
                          subtitle="Каждое правило агента в /case-team имеет явную ссылку на свой уровень источника. Это снижает галлюцинации в стратегических цифрах.")


def build_appendix_a8_quality(prs, slide_num):
    s = add_blank(prs)
    build_appendix_header(s, slide_num, 8, 9,
        "4 параллельных проверки качества дизайна",
        "Перед сдачей презентации case-designer запускает 4 параллельных под-скилла из claude-design — каждый ловит свой класс багов.")

    # 4 карточки 2×2
    checks = [
        ("AI-slop", "ai-slop-check.md",
         "Шаблонные формулировки, generic-фразы, эмодзи как декор, градиенты везде.",
         PRIMARY),
        ("Accessibility", "accessibility-audit.md",
         "Контраст 4.5:1, читаемость заголовков, цвет как единственный носитель смысла.",
         ACCENT_POS),
        ("Hierarchy / Rhythm", "hierarchy-rhythm-review.md",
         "Один visual anchor per slide, ≤7±2 блоков, нет полотен текста.",
         RGBColor(0x6F, 0x46, 0xC0)),
        ("Polish", "polish-pass.md",
         "Финальный гейт перед сдачей.",
         ACCENT_WARN),
    ]
    card_w = Inches(6.10)
    card_h = Inches(1.55)
    gap = Inches(0.20)
    x_l = Inches(0.55)
    x_r = x_l + card_w + gap
    y_top = Inches(2.95)
    y_bot = y_top + card_h + gap
    coords = [(x_l, y_top), (x_r, y_top), (x_l, y_bot), (x_r, y_bot)]

    for (name, skill, body, color), (x, y) in zip(checks, coords):
        add_rect(s, x, y, card_w, card_h, fill=PRIMARY_SOFT, radius=0.06)
        add_rect(s, x, y, Inches(0.10), card_h, fill=color)
        add_textbox(s, x + Inches(0.30), y + Inches(0.15),
                    card_w - Inches(0.50), Inches(0.35),
                    text=name, size=16, bold=True, color=PRIMARY_DARK,
                    font=FONT_DISPLAY)
        add_textbox(s, x + Inches(0.30), y + Inches(0.50),
                    card_w - Inches(0.50), Inches(0.30),
                    text=skill, size=10, color=INK_2, font=FONT_MONO)
        add_textbox(s, x + Inches(0.30), y + Inches(0.82),
                    card_w - Inches(0.50), card_h - Inches(0.90),
                    text=body, size=11, color=INK, line_spacing=1.30)

    # Доп. снизу
    extra_y = y_bot + card_h + Inches(0.25)
    add_rect(s, Inches(0.55), extra_y, SLIDE_W - Inches(1.10), Inches(0.60),
             fill=PRIMARY_DARK, radius=0.06)
    add_textbox(s, Inches(0.75), extra_y + Inches(0.08),
                SLIDE_W - Inches(1.50), Inches(0.22),
                text="ДОПОЛНИТЕЛЬНО", size=10, bold=True, color=PRIMARY,
                font=FONT_DISPLAY)
    add_textbox(s, Inches(0.75), extra_y + Inches(0.28),
                SLIDE_W - Inches(1.50), Inches(0.30),
                text="Pre-flight overflow check (font × lines × spacing / 72 ≤ container_h)  ·  macOS PowerPoint blessing (открыть → Восстановить → Save As)",
                size=11, color=WHITE, line_spacing=1.25)

    draw_footer(s, slide_num)


def build_appendix_a9_business(prs, slide_num):
    s = add_blank(prs)
    build_appendix_header(s, slide_num, 9, 9,
        "Бизнес-модель (открытая)",
        "/case-team — open-source skill для Claude Code. Категориальный сдвиг — это не SaaS-стартап и не «free pilot до цены», а реально бесплатный инструмент для пользователей Claude.")

    # 3 колонки: что есть / это НЕ / монетизация автора
    y0 = Inches(2.95)
    col_w = Inches(4.05)
    col_gap = Inches(0.10)
    col_h = Inches(3.85)

    # Колонка 1
    x1 = Inches(0.55)
    add_rect(s, x1, y0, col_w, col_h, fill=PRIMARY_SOFT, radius=0.06)
    add_textbox(s, x1 + Inches(0.20), y0 + Inches(0.15),
                col_w - Inches(0.40), Inches(0.32),
                text="ЧТО ЕСТЬ /case-team", size=10, bold=True,
                color=PRIMARY_DARK, font=FONT_DISPLAY)
    items = [
        "Open-source skill для Claude Code",
        "Устанавливается одной командой в свою Claude-среду",
        "Стоимость для пользователя — $0 дополнительно, если есть Claude подписка ($20-100/мес уже оплачена)",
    ]
    tb = add_textbox(s, x1 + Inches(0.20), y0 + Inches(0.55),
                     col_w - Inches(0.40), col_h - Inches(0.70),
                     text="", size=12, color=INK, line_spacing=1.40)
    tf = tb.text_frame
    set_text(tf, "•  " + items[0], size=12, color=INK, line_spacing=1.40)
    for it in items[1:]:
        add_paragraph(tf, "•  " + it, size=12, color=INK, line_spacing=1.40,
                      space_before=6)

    # Колонка 2 — это НЕ
    x2 = x1 + col_w + col_gap
    add_rect(s, x2, y0, col_w, col_h, fill=GREY_PALE, radius=0.06)
    add_textbox(s, x2 + Inches(0.20), y0 + Inches(0.15),
                col_w - Inches(0.40), Inches(0.32),
                text="ЭТО НЕ:", size=10, bold=True,
                color=PRIMARY_DARK, font=FONT_DISPLAY)
    nots = [
        "SaaS-стартап (нет отдельной подписки на продукт)",
        "«Free pilot до взлёта цены» — реально бесплатный инструмент",
        "Закрытая корпоративная разработка — open-source MIT",
    ]
    tb = add_textbox(s, x2 + Inches(0.20), y0 + Inches(0.55),
                     col_w - Inches(0.40), col_h - Inches(0.70),
                     text="", size=12, color=INK, line_spacing=1.40)
    tf = tb.text_frame
    set_text(tf, "✕  " + nots[0], size=12, color=INK, line_spacing=1.40)
    for it in nots[1:]:
        add_paragraph(tf, "✕  " + it, size=12, color=INK, line_spacing=1.40,
                      space_before=6)

    # Колонка 3 — Монетизация автора
    x3 = x2 + col_w + col_gap
    add_rect(s, x3, y0, col_w, col_h, fill=PRIMARY_DARK, radius=0.06)
    add_textbox(s, x3 + Inches(0.20), y0 + Inches(0.15),
                col_w - Inches(0.40), Inches(0.32),
                text="МОНЕТИЗАЦИЯ АВТОРА", size=10, bold=True,
                color=PRIMARY, font=FONT_DISPLAY)
    mons = [
        "Узнаваемость в продуктовом сообществе",
        "Возможные консультации",
        "Карьерные возможности",
    ]
    tb = add_textbox(s, x3 + Inches(0.20), y0 + Inches(0.55),
                     col_w - Inches(0.40), col_h - Inches(0.70),
                     text="", size=12, color=WHITE, line_spacing=1.40)
    tf = tb.text_frame
    set_text(tf, "•  " + mons[0], size=12, color=WHITE, line_spacing=1.40)
    for it in mons[1:]:
        add_paragraph(tf, "•  " + it, size=12, color=WHITE, line_spacing=1.40,
                      space_before=6)

    # Footer note
    note_y = Inches(6.95)
    add_textbox(s, Inches(0.55), note_y, SLIDE_W - Inches(1.10), Inches(0.32),
                text="В БУДУЩЕМ: после интеграции Claude Design в десктоп-приложение Anthropic — premium-режим (HTML) тоже автоматизируется, не потребует ручной загрузки YAML в браузер.",
                size=10, italic=True, color=INK_2, line_spacing=1.25)

    draw_footer(s, slide_num)


# ────────────────────────────────────────────────────────────────────────
# MAIN
# ────────────────────────────────────────────────────────────────────────

def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Main content (9 slides)
    build_slide_01_cover(prs)
    build_slide_02_exec(prs)
    build_slide_03_problem(prs)
    build_slide_04_concept(prs)
    build_slide_05_realization(prs)
    build_slide_06_economics(prs)
    build_slide_07_changes(prs)
    build_slide_08_team(prs)
    build_slide_09_final(prs)

    # Appendices (A0 = 3 slides if 71 sources × ~24 per slide)
    build_appendix_a0_sources(prs)            # slides 10-12 (A0.1, A0.2, A0.3)
    # Update appendix numbering — A0 took 3 slots, next slot 13
    next_slide = 10 + 3  # будет 13, если 3 чанка
    build_appendix_a1_segment(prs, next_slide); next_slide += 1
    build_appendix_a2_competitors(prs, next_slide); next_slide += 1
    build_appendix_a3_tech(prs, next_slide); next_slide += 1
    build_appendix_a4_finmodel(prs, next_slide); next_slide += 1
    build_appendix_a5_cascade(prs, next_slide); next_slide += 1
    build_appendix_a6_roadmap(prs, next_slide); next_slide += 1
    build_appendix_a7_methods(prs, next_slide); next_slide += 1
    build_appendix_a8_quality(prs, next_slide); next_slide += 1
    build_appendix_a9_business(prs, next_slide); next_slide += 1

    prs.save(OUT_FILE)
    print(f"✓ saved: {OUT_FILE}")
    return OUT_FILE


# ────────────────────────────────────────────────────────────────────────
# PRE-FLIGHT OVERFLOW CHECK
# ────────────────────────────────────────────────────────────────────────

def preflight_check(pptx_path):
    EMU = 914400
    SLIDE_H_EMU = int(7.5 * EMU)
    issues = []
    with zipfile.ZipFile(pptx_path) as z:
        for name in z.namelist():
            if not (name.startswith("ppt/slides/slide") and name.endswith(".xml")):
                continue
            xml = z.read(name).decode("utf-8")
            for sp in re.findall(r'<p:sp>(.*?)</p:sp>', xml, re.DOTALL):
                off = re.search(r'<a:off x="-?\d+" y="(-?\d+)"\s*/>', sp)
                ext = re.search(r'<a:ext cx="\d+" cy="(\d+)"\s*/>', sp)
                if not (off and ext):
                    continue
                y, cy = int(off.group(1)), int(ext.group(1))
                if y + cy > SLIDE_H_EMU + 2000:
                    text = re.search(r'<a:t>([^<]+)</a:t>', sp)
                    t = text.group(1)[:60] if text else ""
                    issues.append(
                        f"{name}: y+h={(y+cy)/EMU:.2f}\" → '{t}'"
                    )
    if issues:
        print("⚠ PRE-FLIGHT FAILED:")
        for i in issues:
            print(" ", i)
        return False
    print("✓ pre-flight ok — все shapes в пределах слайда (y+h ≤ 7.5\")")
    return True


if __name__ == "__main__":
    out = build()
    preflight_check(out)
